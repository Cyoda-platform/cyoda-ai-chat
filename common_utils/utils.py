import io
import os
import logging
import time

import requests
from typing import Optional
import uuid
import json
import jsonschema
from bs4 import BeautifulSoup
from jsonschema import validate
import tempfile
import magic
import csv
import chardet

import cairosvg


from langchain_community.document_loaders.pdf import PyPDFLoader
from langchain_community.document_loaders.word_document import UnstructuredWordDocumentLoader
from langchain_community.document_loaders import UnstructuredExcelLoader
from langchain_community.document_loaders import CSVLoader

from odf.opendocument import load
from odf.text import P
from odf.table import Table, TableRow, TableCell
from odf import teletype

from nltk import download as nltk_download
from pptx import Presentation

logger = logging.getLogger('django')

def get_user_answer(response):
    answer = response.get('message', '') if response and isinstance(response, dict) else ''
    if isinstance(answer, dict) or isinstance(answer, list):
        answer = json.dumps(answer)
    return answer

def generate_uuid() -> uuid:
    return uuid.uuid1()

def parse_json(result: str) -> str:
    if result.startswith("```"):
        return "\n".join(result.split("\n")[1:-1])
    if not result.startswith("{"):
        start_index = result.find("```json")
        if start_index != -1:
            start_index += len("```json\n")
            end_index = result.find("```", start_index)
            return result[start_index:end_index].strip()
    return result

def validate_result(parsed_result: str, file_path: str) -> bool:
    try:
        with open(file_path, "r") as schema_file:
            schema = json.load(schema_file)
    except (FileNotFoundError, json.JSONDecodeError) as e:
        logger.error(f"Error reading schema file {file_path}: {e}")
        raise

    try:
        json_data = json.loads(parsed_result)
        validator = jsonschema.Draft7Validator(schema)
        errors = sorted(validator.iter_errors(json_data), key=lambda e: e.path)

        if errors:
            error_messages = []
            for error in errors:
                error_messages.append(
                    f"Message: {error.message}\n"
                    f"Schema path: {'.'.join(map(str, error.schema_path))}\n"
                    f"Instance path: {'.'.join(map(str, error.path))}"
                )
            full_error_message = "\n\n".join(error_messages)
            logger.error(f"JSON validation failed:\n{full_error_message}")
            raise jsonschema.exceptions.ValidationError(full_error_message)

        logger.info("JSON validation successful.")
        return True

    except json.JSONDecodeError as e:
        logger.error(f"Failed to decode JSON: {e}")
        raise

def get_env_var(name: str, default_value: Optional[str] = None) -> str:
    value = os.getenv(name)
    if value is None:
        if default_value is None:
            logger.warning(f"Environment variable {name} not found.")
        value = default_value
    return value

def read_file(file_path: str):
    """Read and return JSON data from a file."""
    try:
        with open(file_path, 'r') as file:
            return file.read()
    except Exception as e:
        logger.error(f"Failed to read JSON file {file_path}: {e}")
        logger.exception("An exception occurred")
        raise

def read_json_file(file_path: str):
    try:
        with open(file_path, "r") as file:
            data = json.load(file)
        logger.info(f"Successfully read JSON file: {file_path}")
        return data
    except FileNotFoundError as e:
        logger.error(f"File not found: {file_path}")
        raise
    except json.JSONDecodeError as e:
        logger.error(f"JSON decoding failed for file {file_path}: {e}")
        raise

def send_get_request(token: str, api_url: str, path: str) -> Optional[requests.Response]:
    url = f"{api_url}/{path}"
    token = f"Bearer {token}" if not token.startswith('Bearer') else token
    headers = {
        "Content-Type": "application/json",
        "Authorization": f"{token}",
    }
    try:
        response = requests.get(url, headers=headers)
        #todo response.raise_for_status()  # Raise an error for bad status codes
        logger.info(f"GET request to {url} successful.")
        return response
    except requests.exceptions.HTTPError as http_err:
        logger.error(f"HTTP error during GET request to {url}: {http_err}")
        raise
    except Exception as err:
        logger.error(f"Error during GET request to {url}: {err}")
        logger.exception("An exception occurred")
        raise

def send_post_request(token: str, api_url: str, path: str, data=None, json=None) -> Optional[requests.Response]:
    url = f"{api_url}/{path}"
    token = f"Bearer {token}" if not token.startswith('Bearer') else token
    headers = {
        "Content-Type": "application/json",
        "Authorization": f"{token}",
    }
    try:
        response = requests.post(url, headers=headers, data=data, json=json)
        response.raise_for_status()  # Raise an error for bad status codes
        logger.info(f"POST request to {url} successful.")
        return response
    except requests.exceptions.HTTPError as http_err:
        logger.error(f"HTTP error during POST request to {url}: {http_err}")
        raise
    except Exception as err:
        logger.error(f"Error during POST request to {url}: {err}")
        logger.exception("An exception occurred")
        raise

def send_put_request(token: str, api_url: str, path: str, data=None, json=None) -> Optional[requests.Response]:
    url = f"{api_url}/{path}"
    token = f"Bearer {token}" if not token.startswith('Bearer') else token
    headers = {
        "Content-Type": "application/json",
        "Authorization": f"{token}",
    }
    try:
        response = requests.put(url, headers=headers, data=data, json=json)
        response.raise_for_status()  # Raise an error for bad status codes
        logger.info(f"PUT request to {url} successful.")
        return response
    except requests.exceptions.HTTPError as http_err:
        logger.error(f"HTTP error during PUT request to {url}: {http_err}")
        raise
    except Exception as err:
        logger.exception("An exception occurred")
        logger.error(f"Error during PUT request to {url}: {err}")
        raise

def send_delete_request(token: str, api_url: str, path: str) -> Optional[requests.Response]:
    url = f"{api_url}/{path}"
    token = f"Bearer {token}" if not token.startswith('Bearer') else token
    headers = {
        "Content-Type": "application/json",
        "Authorization": f"{token}",
    }
    try:
        response = requests.delete(url, headers=headers)
        response.raise_for_status()  # Raise an error for bad status codes
        logger.info(f"GET request to {url} successful.")
        return response
    except requests.exceptions.HTTPError as http_err:
        logger.error(f"HTTP error during GET request to {url}: {http_err}")
        raise
    except Exception as err:
        logger.exception("An exception occurred")
        logger.error(f"Error during GET request to {url}: {err}")
        raise

def expiration_date(seconds: int) -> int:
    return int((time.time()+seconds)*1000.0)

def now():
    timestamp = int(time.time()*1000.0)
    return timestamp

def timestamp_before(seconds: int) -> int:
    return int((time.time()-seconds)*1000.0)

def validate_and_parse_json_v1(
            processor,
            chat_id: str,
            data: str,
            schema_path: str,
            max_retries: int,
    ):
        try:
            parsed_data = parse_json(data)
        except Exception as e:
            logger.exception("An exception occurred")
            logger.error(f"Failed to parse JSON: {e}")
            raise ValueError("Invalid JSON data provided.") from e

        attempt = 0
        while attempt <= max_retries:
            try:
                validate_result(parsed_data, schema_path)
                logger.info(f"JSON validation successful on attempt {attempt + 1}.")
                attempt += 1
                return json.loads(parsed_data)
            except jsonschema.exceptions.ValidationError as e:
                attempt += 1
                logger.warning(
                    f"JSON validation failed on attempt {attempt + 1} with error: {e.message}"
                )
                if attempt < max_retries:
                    question = (
                        f"Retry the last step. JSON validation failed with error: {e.message}. "
                        "Return only the DTO JSON."
                    )
                    retry_result = processor.ask_question(chat_id, question)
                    parsed_data = parse_json(retry_result)
            except Exception as e:
                logger.exception("An exception occurred")
                logger.error("Maximum retry attempts reached. Validation failed.")
            finally:
                attempt += 1
        logger.error("Maximum retry attempts reached. Validation failed.")
        raise ValueError("JSON validation failed after retries.")

def validate_and_parse_json(
        processor,
        chat_id: str,
        data: str,
        schema_path: str,
        max_retries: int,
):
    try:
        parsed_data = parse_json(data)
    except Exception as e:
        logger.exception("An exception occurred")
        logger.error(f"Failed to parse JSON: {e}")
        raise ValueError("Invalid JSON data provided.") from e

    attempt = 0
    while attempt <= max_retries:
        try:
            validate_result(parsed_data, schema_path)
            logger.info(f"JSON validation successful on attempt {attempt + 1}.")
            attempt += 1
            return json.loads(parsed_data)
        except jsonschema.exceptions.ValidationError as e:
            attempt += 1
            logger.warning(
                f"JSON validation failed on attempt {attempt + 1} with error: {e.message}"
            )
            try:
                with open(schema_path, "r") as schema_file:
                    schema = json.load(schema_file)
            except (FileNotFoundError, json.JSONDecodeError) as e:
                logger.error(f"Error reading schema file {schema_path}: {e}")
                raise
            if attempt < max_retries:
                question = (
                    f"Retry the last step. JSON validation failed with error: {e.message}. Correct this data: {parsed_data} "
                    f"using this schema: {json.dumps(schema)}. "
                    "Return only the DTO JSON."
                )
                retry_result = processor.ask_question(chat_id, question)
                parsed_data = parse_json(retry_result)
        except Exception as e:
            logger.exception("An exception occurred")
            logger.error("Maximum retry attempts reached. Validation failed.")
        finally:
            attempt += 1
    logger.error("Maximum retry attempts reached. Validation failed.")
    raise ValueError("JSON validation failed after retries.")

def process_uploaded_file(self, file):
    mime_type = get_mime_type(file)
    file_content, metadata = process_uploaded_file_by_mime(self, file, mime_type)
    if "Unsupported file type" in metadata.get("error", ""):
        extension = get_file_extension(file.name)
        mime_type = guess_mime_type_from_extension(extension)
        if mime_type:
            file_content, metadata = process_uploaded_file_by_mime(self, file, mime_type)
    return file_content, metadata

def process_uploaded_file_by_mime(self, file, mime_type):
    """
    Processes a file based on its MIME type.
    Returns the file's contents and metadata.
    """
    try:
        if mime_type == "image/svg+xml":
            image_bytes = convert_svg_to_png(file)
            image_file = io.BytesIO(image_bytes)
            mime_type = get_mime_type(image_file)
            content = self.processor.image_description(image_file, mime_type)
            return content, {"type": "svg", "filename": file.name}

        elif mime_type.startswith("image/"):
            content = self.processor.image_description(file, mime_type)
            return content, {"type": "image", "filename": file.name}

        elif mime_type == "application/pdf":
            file_path = save_temp_file(file)
            content = process_pdf(file_path)
            return content, {"type": "pdf", "filename": file.name}

        elif mime_type in ["application/msword", "application/vnd.openxmlformats-officedocument.wordprocessingml.document"]:
            file_path = save_temp_file(file)
            content = process_word(file_path)
            return content, {"type": "word", "filename": file.name}

        elif mime_type in ["text/plain", "application/json"]:
            content = process_text_file(file)
            return content, {"type": "text", "filename": file.name}

        elif mime_type == "text/markdown":
            content = process_text_file(file)
            return content, {"type": "markdown", "filename": file.name}

        elif mime_type in ["application/vnd.ms-excel",
                           "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"]:
            file_path = save_temp_file(file)
            content = process_excel(file_path)
            return content, {"type": "excel", "filename": file.name}

        elif mime_type in ["text/csv"]:
            encoding = detect_encoding(file)
            encoding_write = "utf-8" # supported by CSVLoader
            file_path = save_temp_file(file, encoding, encoding_write)
            if not has_header(file_path):
                add_fake_header(file_path)
            content = process_csv(file_path)
            return content, {"type": "csv", "filename": file.name}

        elif mime_type in ["application/vnd.ms-powerpoint",
                           "application/vnd.openxmlformats-officedocument.presentationml.presentation"]:
            content = process_presentation(file)
            return content, {"type": "presentation", "filename": file.name}

        elif mime_type in ["text/x-python", "text/x-script.python", "text/x-java-source", "text/x-java", "application/x-php", "application/javascript"]:
            content = process_text_file(file)
            return content, {"type": "code", "filename": file.name}

        elif mime_type in ["application/xml", "text/html", "text/xml"]:
            content = process_xml_or_html(file)
            return content, {"type": "xml_html", "filename": file.name}

        elif mime_type in ["application/vnd.oasis.opendocument.text"]:
            content = process_odt(file)
            return content, {"type": "odt", "filename": file.name}

        elif mime_type == "application/vnd.oasis.opendocument.spreadsheet":
            content = process_ods(file)
            return content, {"type": "ods", "filename": file.name}

        else:
            return None, {"error": f"Unsupported file type: {mime_type}"}
    except Exception as e:
        return None, {"error": f"Error processing file '{file.name}': {str(e)}"}
    finally:
        if 'file_path' in locals() and os.path.exists(file_path):
            os.remove(file_path)

def get_mime_type(file):
    """
    Specifies the MIME type of the uploaded file using python-magic.
    Scans the contents of the file instead of relying on its extension.
    """
    mime = magic.Magic(mime=True)
    mime_type = mime.from_buffer(file.read(1024))
    file.seek(0)
    return mime_type

def guess_mime_type_from_extension(extension):
    """
    Guesses MIME type based on file extension.
    """
    mime_types_by_extension = {
        ".txt": "text/plain",
        ".csv": "text/csv",
        ".json": "application/json",
        ".xml": "application/xml",
        ".html": "text/html",
        ".pdf": "application/pdf",
        ".doc": "application/msword",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".xls": "application/vnd.ms-excel",
        ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ".ppt": "application/vnd.ms-powerpoint",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".py": "text/x-python",
        ".java": "text/x-java-source",
        ".php": "application/x-php",
        ".js": "application/javascript",
        ".jpg": "image/jpeg",
        ".jpeg": "image/jpeg",
        ".png": "image/png",
        ".gif": "image/gif",
    }
    return mime_types_by_extension.get(extension.lower(), None)

def get_file_extension(filename):
    """
    Extracts the extension from a file name.
    """
    return f".{filename.split('.')[-1].lower()}" if "." in filename else ""

def save_temp_file(file, encoding_read=None, encoding_write=None):
    """
    Saves the downloaded file temporarily to disk and returns the path to the file.
    """
    temp_dir = tempfile.gettempdir()
    if not hasattr(file, 'name'):
        raise ValueError("The file object does not have a 'name' attribute.")
    temp_file_path = os.path.join(temp_dir, file.name)

    if encoding_read:
        encoding_write = encoding_write or encoding_read
        with open(temp_file_path, "w", encoding=encoding_write) as f:
            content = file.read().decode(encoding_read)
            f.write(content)
        return temp_file_path

    with open(temp_file_path, "wb") as f:
        content = file.read()
        f.write(content)
    return temp_file_path

def process_pdf(file_path):
    loader = PyPDFLoader(file_path)
    documents = loader.load()
    return "\n".join([doc.page_content for doc in documents])

def process_word(file_path):
    nltk_download('averaged_perceptron_tagger')
    nltk_download('averaged_perceptron_tagger_eng')
    loader = UnstructuredWordDocumentLoader(file_path)
    documents = loader.load()
    return "\n".join([doc.page_content for doc in documents])

def process_text_file(file):
    encoding = detect_encoding(file)
    file_stream = io.StringIO(file.read().decode(encoding))
    content = file_stream.getvalue()
    return content

def process_excel(file_path):
    nltk_download("punkt_tab")
    loader = UnstructuredExcelLoader(file_path)
    documents = loader.load()
    return "\n".join([doc.page_content for doc in documents])

def process_csv(file_path):
    loader = CSVLoader(file_path)
    documents = loader.load()
    return "\n".join([doc.page_content for doc in documents])

def process_presentation(file):
    prs = Presentation(file)
    content = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                content.append(shape.text)
    return "\n".join(content)

def process_xml_or_html(file):
    encoding = detect_encoding(file)
    file_stream = io.StringIO(file.read().decode(encoding))
    xml_content = file_stream.getvalue()
    return xml_content

def process_odt(file):
    odt_doc = load(file)
    paragraphs = []
    for paragraph in odt_doc.getElementsByType(P):
        paragraphs.append(str(paragraph))
    return "\n".join(paragraphs)

def process_ods(file):
    ods_doc = load(file)
    tables = ods_doc.getElementsByType(Table)
    data = []
    for table in tables:
        table_data = []
        rows = table.getElementsByType(TableRow)
        for row in rows:
            row_data = []
            cells = row.getElementsByType(TableCell)
            for cell in cells:
                # Extract text content from the cell using teletype.extractText
                text = ""
                paragraphs = cell.getElementsByType(P)
                for paragraph in paragraphs:
                    text += teletype.extractText(paragraph).strip() if paragraph else ''
                if text:
                    row_data.append(text)
                else:
                    row_data.append('')  # Append empty string if no text is found
            table_data.append(row_data)
        data.append(table_data)
    return data

def convert_svg_to_png(file):
    """
    Processes SVG files and converts them to PNG format.
    Returns the PNG content as bytes.
    """
    svg_content = clean_svg(file)
    png_content = cairosvg.svg2png(bytestring=svg_content)
    return png_content

def clean_svg(svg_content):
    soup = BeautifulSoup(svg_content, "xml")
    for tag in soup.find_all(True):
        for attr in ["content", "data-cell-id", "style"]:
            if attr in tag.attrs:
                del tag[attr]
    for div in soup.find_all("div"):
        div.decompose()
    return str(soup)

def detect_encoding(file):
    try:
        file.seek(0)
        raw_data = file.read(1024 * 1024)
        result = chardet.detect(raw_data)
        file.seek(0)
        return result['encoding']
    except Exception:
        return "utf-8"  # Default fallback encoding

def has_header(file_path):
    """
    Determines if a CSV file has a header row.
    Returns True if a header is detected, False otherwise.
    """
    with open(file_path, 'r', encoding='utf-8') as f:
        sample = f.read(1024)  # Read a small portion of the file
        f.seek(0)  # Reset pointer to the beginning
        dialect = csv.Sniffer().sniff(sample)  # Detect delimiter
        f.seek(0)
        has_header = csv.Sniffer().has_header(sample)  # Check for header
        return has_header


def add_fake_header(file_path):
    """
    Adds a fake header to a CSV file if no header is detected.
    """
    with open(file_path, 'r', encoding='utf-8') as f:
        content = f.readlines()

    # Generate fake headers based on the first row's number of columns
    fake_header = ",".join([f"Column_{i+1}" for i in range(len(content[0].split(",")))])
    content.insert(0, fake_header + "\n")  # Add fake header at the beginning

    with open(file_path, 'w', encoding='utf-8') as f:
        f.writelines(content)

def append_file_content_to_question(question, file_content, metadata):
    """
    Adds file content and metadata to the question text.
    Args:
    question (str): User text query.
    file_content (str): File content extracted using loaders or image analysis.
    metadata (dict): File metadata (type, name, error).
    Returns:
    str: Updated query text.
    """
    file_info = f"File Type: {metadata['type']}, Filename: {metadata['filename']}"
    if question:
        return f"{question}\n\n[Attached File Info: {file_info}]\n\nExtracted Content:\n{file_content}"
    else:
        return f"[Attached File Info: {file_info}]\n\nExtracted Content:\n{file_content}"
